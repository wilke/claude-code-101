#!/usr/bin/env python3
"""Build slides-argonne.pptx from the docs/slides/ markdown ground truth,
on top of the Argonne 16x9 PowerPoint template.

Usage:
    python scripts/build_pptx.py [--template PATH] [--out PATH]

The markdown in docs/slides/ is the single source of truth (see
docs/slides/INDEX.md for the slide format). This script parses it into a block
model and renders each slide onto the template's native layouts:

    layout  -> template layout name
    ------     --------------------
    title   -> "Title Slide"              (aerial-photo background)
    section -> "*Section Break"           (part / exercise dividers)
    content -> "*Title and Subtitle Only" (white base slide, blue left bar)
    closing -> "Closing slide Argonne"

Requires: python-pptx (`pip install python-pptx`).
"""
from __future__ import annotations

import argparse
import copy
import os
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:  # pragma: no cover
    sys.exit("python-pptx is required:  pip install python-pptx")

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slides_md import (  # noqa: E402  shared markdown parser (single source of truth)
    LINK_RE, TOKEN_RE, strip_inline, load_all_slides,
)

REPO = Path(__file__).resolve().parent.parent
SLIDES_DIR = REPO / "docs" / "slides"
DEFAULT_TEMPLATE_CANDIDATES = [
    os.environ.get("ARGONNE_TEMPLATE", ""),
    str(REPO / "templates" / "Argonne_16x9_template.potx"),
    str(Path.home() / "Downloads" / "Argonne_16x9 Presentation Template.potx"),
]

# ---- Argonne palette -------------------------------------------------------
BLUE = RGBColor(0x00, 0x82, 0xCA)        # primary Argonne blue
DARKBLUE = RGBColor(0x00, 0x60, 0x9C)    # accent2 (left bar / subtitle strip)
GOLD = RGBColor(0xF8, 0xB2, 0x00)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x5B, 0x5B, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF2, 0xF1, 0xEC)
CODE_INK = RGBColor(0x1C, 0x1C, 0x1C)
CALLOUT_BG = RGBColor(0xE7, 0xEE, 0xF5)
TABLE_HDR = RGBColor(0x00, 0x60, 0x9C)
TABLE_ALT = RGBColor(0xF4, 0xF7, 0xFA)

MONO = "Consolas"
SANS = "Arial"

# Content region on a content slide (below title + blue kicker strip).
BODY_LEFT = Inches(0.67)
BODY_TOP = Inches(2.3)
BODY_WIDTH = Inches(12.0)
BODY_BOTTOM = Inches(6.88)


# ---------------------------------------------------------------------------
# Inline formatting -> runs  (markdown parsing lives in slides_md.py)
# ---------------------------------------------------------------------------
def add_runs(paragraph, text: str, base_size: int, color=INK,
             font=SANS, bold_default=False):
    """Add richly-formatted runs, honoring **bold** and `code` and [links]."""
    text = LINK_RE.sub(r"\1", text)
    pieces = TOKEN_RE.split(text)
    for piece in pieces:
        if not piece:
            continue
        run = paragraph.add_run()
        if piece.startswith("**") and piece.endswith("**"):
            run.text = piece[2:-2]
            run.font.bold = True
        elif piece.startswith("`") and piece.endswith("`"):
            run.text = piece[1:-1]
            run.font.name = MONO
        else:
            # collapse leftover single-underscore emphasis
            run.text = re.sub(r"(?<!\w)_([^_]+)_(?!\w)", r"\1", piece)
            run.font.bold = bold_default
            run.font.name = font
        run.font.size = Pt(base_size)
        run.font.color.rgb = color
        if run.font.name is None:
            run.font.name = font


# ---------------------------------------------------------------------------
# Template handling
# ---------------------------------------------------------------------------
def potx_to_pptx(potx_path: Path) -> Path:
    """Return a .pptx copy of a .potx (python-pptx rejects template content type)."""
    if potx_path.suffix.lower() == ".pptx":
        return potx_path
    tmp = Path(tempfile.mkdtemp(prefix="argonne_tpl_")) / "template.pptx"
    with zipfile.ZipFile(potx_path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            data = zin.read(item)
            if item == "[Content_Types].xml":
                data = data.replace(
                    b"presentationml.template.main+xml",
                    b"presentationml.presentation.main+xml",
                )
            zout.writestr(item, data)
    return tmp


def delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def layouts_by_name(prs) -> dict:
    out = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            out[layout.name] = layout
    return out


def get_ph(slide, idx=None, ph_type=None):
    for ph in slide.placeholders:
        fmt = ph.placeholder_format
        if idx is not None and fmt.idx == idx:
            return ph
        if ph_type is not None and fmt.type == ph_type:
            return ph
    return None


def set_frame_defaults(tf, wrap=True):
    tf.word_wrap = wrap
    try:
        tf.auto_size = None
    except Exception:
        pass
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


# ---------------------------------------------------------------------------
# Height estimation (very rough, in inches) so blocks stack without huge gaps
# ---------------------------------------------------------------------------
def wrapped_lines(text: str, chars_per_line: int) -> int:
    return max(1, -(-len(strip_inline(text)) // chars_per_line))


def est_height(block, width_in: float) -> float:
    cpl = int(width_in / 0.105)  # ~chars per line at ~15pt Arial (slight overestimate)
    if block["type"] == "lede":  # 17pt — fewer chars per line, taller
        lcpl = int(width_in / 0.125)
        return 0.06 + wrapped_lines(block["text"], lcpl) * 0.34 + 0.12
    if block["type"] == "para":
        return 0.04 + wrapped_lines(block["text"], cpl) * 0.30 + 0.08
    if block["type"] == "bullets":
        h = 0.06
        for it in block["items"]:
            h += wrapped_lines(it["text"], cpl - 4) * 0.31 + 0.06
        return h
    if block["type"] == "numbered":
        h = 0.06
        for it in block["items"]:
            h += wrapped_lines(it, cpl - 4) * 0.31 + 0.07
        return h
    if block["type"] == "code":
        nlines = len(block["lines"])
        fs = code_font_size(nlines)
        return nlines * (fs / 72.0 * 1.28) + 0.28
    if block["type"] == "table":
        rows = block["rows"]
        ncols = max(len(r) for r in rows)
        colw = width_in / ncols
        ccpl = max(6, int(colw / 0.085))  # ~chars/col-line at 11.5pt
        h = 0.06
        for r in rows:
            maxlines = max((wrapped_lines(c, ccpl) for c in r), default=1)
            h += maxlines * 0.26 + 0.16
        return h
    if block["type"] == "callout":
        return 0.24 + wrapped_lines(block["text"], cpl - 2) * 0.28
    if block["type"] == "columns":
        maxi = max((len(c["items"]) for c in block["cols"]), default=0)
        return 0.4 + maxi * 0.42
    return 0.4


def code_font_size(nlines: int) -> int:
    if nlines > 24:
        return 9
    if nlines > 18:
        return 10
    if nlines > 12:
        return 11
    return 12


# ---------------------------------------------------------------------------
# Renderers
# ---------------------------------------------------------------------------
def remove_placeholder(slide, idx=None, ph_type=None):
    ph = get_ph(slide, idx=idx, ph_type=ph_type)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def render_title(slide, data, kicker):
    # drop the empty picture placeholder (renders as a gray box)
    from pptx.enum.shapes import PP_PLACEHOLDER
    remove_placeholder(slide, ph_type=PP_PLACEHOLDER.PICTURE)
    remove_placeholder(slide, idx=10)
    # ctrTitle
    ct = slide.shapes.title
    if ct is not None:
        ct.text_frame.clear()
        p = ct.text_frame.paragraphs[0]
        add_runs(p, data["title"], 44, color=WHITE, bold_default=True)
    subtitle = get_ph(slide, idx=1)
    if subtitle is not None:
        subtitle.text_frame.clear()
        p = subtitle.text_frame.paragraphs[0]
        add_runs(p, kicker.upper() if kicker else "", 14, color=GOLD, bold_default=True)
    # description + contributors -> manual white textbox lower-left
    lines = []
    if data["lede"]:
        lines.append(("lede", data["lede"]))
    for b in data["blocks"]:
        if b["type"] == "para":
            lines.append(("para", b["text"]))
    box = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(6.7), Inches(2.7))
    tf = box.text_frame
    set_frame_defaults(tf)
    first = True
    for kind, text in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if text.startswith("Contributors:"):
            p.space_before = Pt(10)
            lbl, names = text.split(":", 1)
            r = p.add_run(); r.text = "CONTRIBUTORS  "
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = SANS
            r2 = p.add_run(); r2.text = names.strip()
            r2.font.size = Pt(13); r2.font.color.rgb = WHITE; r2.font.name = SANS
        elif kind == "lede":
            add_runs(p, text, 17, color=WHITE, bold_default=True)
            p.space_after = Pt(8)
        else:
            add_runs(p, text, 13, color=RGBColor(0xEA, 0xF2, 0xF8))
            p.space_after = Pt(6)


def render_section(slide, data, kicker, links_blocks, capstones=None):
    title_ph = slide.shapes.title
    # standardize the box (keep the template's tall, middle-anchored geometry) so
    # long ledes wrap comfortably inside the slide instead of hugging the edge.
    # When a capstone list follows, anchor the title/lede to the top instead.
    title_ph.left = Inches(0.6)
    title_ph.width = Inches(11.9)
    if capstones:
        title_ph.top = Inches(0.35)
        title_ph.height = Inches(1.9)
    else:
        title_ph.top = Inches(0)
        title_ph.height = Inches(6.55)
    tf = title_ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_right = Inches(0.1)
    if capstones:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    # kicker
    p0 = tf.paragraphs[0]
    if kicker:
        r = p0.add_run(); r.text = kicker.upper()
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = SANS
        pt = tf.add_paragraph()
    else:
        pt = p0
    add_runs(pt, data["title"], 46, color=WHITE, bold_default=True)
    pt.space_before = Pt(6)
    if data["lede"]:
        pl = tf.add_paragraph()
        add_runs(pl, data["lede"], 18, color=RGBColor(0xEA, 0xF2, 0xF8))
        pl.space_before = Pt(10)

    # exercise-divider link columns (rendered as manual textboxes)
    if links_blocks:
        cols = links_blocks["cols"]
        ncol = len(cols)
        total_w = 12.0
        gap = 0.4
        cw = (total_w - gap * (ncol - 1)) / ncol
        x = 0.7
        for c in cols:
            box = slide.shapes.add_textbox(Inches(x), Inches(4.6), Inches(cw), Inches(2.2))
            bt = box.text_frame
            set_frame_defaults(bt)
            hp = bt.paragraphs[0]
            r = hp.add_run(); r.text = c["heading"].upper()
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = SANS
            for it in c["items"]:
                ip = bt.add_paragraph()
                add_runs(ip, it, 12.5, color=WHITE, font=MONO)
                ip.space_before = Pt(4)
            x += cw + gap

    # capstone showcase — static "title | description" columns, row-aligned.
    # The HTML deck lists the titles and rotates the description instead; here every
    # description is one line, so two matched-paragraph textboxes stay aligned.
    if capstones:
        items = capstones["items"]
        light = RGBColor(0xEA, 0xF2, 0xF8)
        cols = [
            {"x": 0.7, "w": 3.3, "bold": True, "color": WHITE, "key": "name"},
            {"x": 4.1, "w": 8.5, "bold": False, "color": light, "key": "desc"},
        ]
        for col in cols:
            box = slide.shapes.add_textbox(Inches(col["x"]), Inches(2.5),
                                           Inches(col["w"]), Inches(4.2))
            bt = box.text_frame
            set_frame_defaults(bt)
            bt.word_wrap = True
            first = True
            for it in items:
                p = bt.paragraphs[0] if first else bt.add_paragraph()
                first = False
                p.space_after = Pt(10)
                r = p.add_run(); r.text = it.get(col["key"], "")
                r.font.size = Pt(14); r.font.bold = col["bold"]
                r.font.color.rgb = col["color"]; r.font.name = SANS


def render_closing(slide, data, kicker):
    # Closing layout has decorative freeforms; drop any inherited placeholders
    # (they render as stray prompt glyphs) and add a centered white title+lede box.
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, data["title"], 54, color=WHITE, bold_default=True)
    if data["lede"]:
        pl = tf.add_paragraph()
        add_runs(pl, data["lede"], 20, color=RGBColor(0xEA, 0xF2, 0xF8))
        pl.space_before = Pt(14)


def render_content(slide, data, kicker):
    # title
    title_ph = slide.shapes.title
    if title_ph is not None:
        title_ph.text_frame.clear()
        p = title_ph.text_frame.paragraphs[0]
        add_runs(p, data["title"], 30, color=DARKBLUE, bold_default=True)
    # blue kicker strip (idx 13) — set fill explicitly (layout fill doesn't inherit reliably)
    strip = get_ph(slide, idx=13)
    if strip is not None:
        strip.fill.solid()
        strip.fill.fore_color.rgb = DARKBLUE
        strip.line.fill.background()
        tf = strip.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        sp = tf.paragraphs[0]
        r = sp.add_run(); r.text = (kicker or "").upper()
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = SANS

    blocks = list(data["blocks"])
    if data["lede"]:
        blocks.insert(0, {"type": "lede", "text": data["lede"]})

    avail = (BODY_BOTTOM - BODY_TOP) / 914400  # inches
    total = sum(est_height(b, 11.9) for b in blocks)
    # Scale FONTS and positions together on dense slides so text shrinks with
    # the boxes (prevents overlap). Floor keeps text readable; below it we allow
    # mild bottom overflow rather than crushing the type.
    fs = min(1.0, avail / total) if total > 0 else 1.0
    fs = max(fs, 0.74)

    y = BODY_TOP / 914400
    for b in blocks:
        h = est_height(b, 11.9) * fs
        render_block(slide, b, y, h, fs)
        y += h + 0.07 * fs


def render_block(slide, b, y, h, fs=1.0):
    left = BODY_LEFT
    width = BODY_WIDTH
    t = b["type"]

    def S(pt):
        return pt * fs

    if t in ("para", "lede"):
        box = slide.shapes.add_textbox(left, Inches(y), width, Inches(max(h, 0.3)))
        tf = box.text_frame; set_frame_defaults(tf)
        p = tf.paragraphs[0]
        if t == "lede":
            add_runs(p, b["text"], S(17), color=MUTED)
        else:
            add_runs(p, b["text"], S(15), color=INK)

    elif t == "bullets":
        box = slide.shapes.add_textbox(left, Inches(y), width, Inches(max(h, 0.3)))
        tf = box.text_frame; set_frame_defaults(tf)
        first = True
        for it in b["items"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = it["level"]
            bullet = "– " if it["level"] else "•  "
            r = p.add_run(); r.text = bullet
            r.font.size = Pt(S(15)); r.font.color.rgb = BLUE; r.font.name = SANS
            add_runs(p, it["text"], S(15), color=INK)
            p.space_after = Pt(S(3))

    elif t == "numbered":
        box = slide.shapes.add_textbox(left, Inches(y), width, Inches(max(h, 0.3)))
        tf = box.text_frame; set_frame_defaults(tf)
        first = True
        for k, it in enumerate(b["items"], 1):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = f"{k}.  "
            r.font.size = Pt(S(15)); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = SANS
            add_runs(p, it, S(15), color=INK)
            p.space_after = Pt(S(3))

    elif t == "code":
        cfs = code_font_size(len(b["lines"])) * fs
        box = slide.shapes.add_textbox(left, Inches(y), width, Inches(max(h, 0.3)))
        box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
        box.line.color.rgb = RGBColor(0xDD, 0xD8, 0xCC); box.line.width = Pt(0.75)
        # accent left edge
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(y),
                                     Pt(3), Inches(max(h, 0.3)))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
        tf = box.text_frame; set_frame_defaults(tf)
        tf.margin_left = Inches(0.12)
        first = True
        for ln in b["lines"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = ln if ln else " "
            r.font.name = MONO; r.font.size = Pt(cfs); r.font.color.rgb = CODE_INK
            p.line_spacing = 1.05

    elif t == "callout":
        box = slide.shapes.add_textbox(left, Inches(y), width, Inches(max(h, 0.35)))
        box.fill.solid(); box.fill.fore_color.rgb = CALLOUT_BG
        box.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(y),
                                     Pt(3), Inches(max(h, 0.35)))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
        tf = box.text_frame; set_frame_defaults(tf)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        p = tf.paragraphs[0]
        add_runs(p, b["text"], S(14), color=RGBColor(0x14, 0x3A, 0x5A))

    elif t == "table":
        render_table(slide, b["rows"], left, Inches(y), width, Inches(max(h, 0.4)), fs)

    elif t == "columns":
        cols = b["cols"]
        ncol = len(cols)
        gap = 0.35
        cw = (width / 914400 - gap * (ncol - 1)) / ncol
        x = left / 914400
        for c in cols:
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cw), Inches(max(h, 0.4)))
            tf = box.text_frame; set_frame_defaults(tf)
            hp = tf.paragraphs[0]
            r = hp.add_run(); r.text = c["heading"].upper()
            r.font.size = Pt(S(11.5)); r.font.bold = True; r.font.color.rgb = DARKBLUE; r.font.name = SANS
            hp.space_after = Pt(S(4))
            for it in c["items"]:
                ip = tf.add_paragraph()
                is_link = it.strip().startswith("[")
                r = ip.add_run(); r.text = "•  "
                r.font.size = Pt(S(13)); r.font.color.rgb = BLUE; r.font.name = SANS
                add_runs(ip, it, S(13), color=INK, font=(MONO if is_link else SANS))
                ip.space_after = Pt(S(2))
            x += cw + gap


def render_table(slide, rows, left, top, width, height, fs=1.0):
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    gtbl = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    tbl.first_row = True
    tbl.horz_banding = False
    for ci in range(ncols):
        for ri in range(nrows):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            text = rows[ri][ci] if ci < len(rows[ri]) else ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HDR
                add_runs(p, text, 11.5 * fs, color=WHITE, bold_default=True)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT if ri % 2 == 0 else WHITE
                add_runs(p, text, 11.5 * fs, color=INK)


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default=None)
    ap.add_argument("--out", default=str(REPO / "slides-argonne.pptx"))
    args = ap.parse_args()

    template = args.template
    if not template:
        for cand in DEFAULT_TEMPLATE_CANDIDATES:
            if cand and Path(cand).exists():
                template = cand
                break
    if not template or not Path(template).exists():
        sys.exit("Argonne template not found. Pass --template PATH or set "
                 "ARGONNE_TEMPLATE / place it at templates/Argonne_16x9_template.potx")

    pptx_template = potx_to_pptx(Path(template))
    prs = Presentation(str(pptx_template))
    delete_all_slides(prs)
    L = layouts_by_name(prs)

    section_layout = L.get("*Section Break")
    layout_map = {
        "title": L.get("Title Slide"),
        "section": section_layout,
        "content": L.get("*Title and Subtitle Only"),
        # The capstone reads as a section divider in the source deck; render it on
        # the aerial "Section Break" layout for visual consistency (the dedicated
        # "Closing" layout carries decorative freeforms that render as stray slivers).
        "closing": section_layout,
    }

    slides = load_all_slides(SLIDES_DIR)
    for attrs, data in slides:
        kind = attrs.get("layout", "content")
        kicker = attrs.get("kicker", "")
        layout = layout_map.get(kind) or layout_map["content"]
        slide = prs.slides.add_slide(layout)

        # pull out a links-columns block (exercise dividers) / capstones block
        links = capstones = None
        if kind in ("section", "closing"):
            for b in data["blocks"]:
                if b["type"] == "columns" and links is None:
                    links = b
                elif b["type"] == "capstones" and capstones is None:
                    capstones = b

        if kind == "title":
            render_title(slide, data, kicker)
        elif kind in ("section", "closing"):
            render_section(slide, data, kicker, links, capstones)
        else:
            render_content(slide, data, kicker)

    prs.save(args.out)
    print(f"Wrote {args.out}  ({len(slides)} slides)")


if __name__ == "__main__":
    main()
